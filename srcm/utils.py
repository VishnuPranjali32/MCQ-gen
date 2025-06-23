import re
import fitz  # PyMuPDF for PDFs
from docx import Document
from pptx import Presentation
import tempfile


def clean_text(text: str) -> str:
    """Clean input text by removing extra whitespace and special characters."""
    text = re.sub(r'\s+', ' ', text.strip())
    return text

def validate_inputs(input_text: str, num_questions: int, num_options: int, difficulty_config) -> None:
    """Validate user inputs."""
    if not input_text:
        raise ValueError("Input text or topic cannot be empty.")
    if num_questions < 1:
        raise ValueError("Number of questions must be at least 1.")
    if num_options < 2:
        raise ValueError("Number of options must be at least 2.")

    allowed_levels = {"easy", "medium", "hard"}

    if isinstance(difficulty_config, dict):
        total = 0
        for level, count in difficulty_config.items():
            if level not in allowed_levels:
                raise ValueError(f"Invalid difficulty level: {level}")
            if not isinstance(count, int) or count < 0:
                raise ValueError(f"Invalid count for difficulty '{level}': must be a non-negative integer")
            total += count
        if total != num_questions:
            raise ValueError("Sum of custom difficulty question counts must equal total number of questions.")
    elif isinstance(difficulty_config, str):
        if difficulty_config.lower() != "mixed":
            raise ValueError("Difficulty must be 'mixed' or a dictionary of counts.")
    else:
        raise ValueError("Difficulty configuration must be a string or dictionary.")


def extract_text_from_file(uploaded_file) -> str:
    """Extract and clean text from PDF, DOCX, or PPTX Streamlit upload."""
    suffix = uploaded_file.name.split(".")[-1].lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{suffix}") as temp_file:
        temp_file.write(uploaded_file.read())
        temp_path = temp_file.name

    text = ""

    if suffix == "pdf":
        doc = fitz.open(temp_path)
        text = "\n".join(page.get_text() for page in doc)
    elif suffix == "docx":
        doc = Document(temp_path)
        text = "\n".join(p.text for p in doc.paragraphs)
    elif suffix == "pptx":
        prs = Presentation(temp_path)
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    else:
        raise ValueError("Unsupported file type. Please upload PDF, DOCX, or PPTX.")

    return clean_text(text)
